from concurrent.futures import thread
import os
from turtle import width
from pikepdf import Pdf
from tkinter import filedialog
from tkinter import *
from tkinter import simpledialog

import img2pdf
import PIL.Image
import io
from pdf2image import convert_from_path
 

# pdf -> ppt
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches

from io import BytesIO


def mergeBT():
    # fs = str(filedialog.askdirectory()).replace("/","\\")

    lblMergeBTN.config(text="Working!")
    files = filedialog.askopenfiles()
    p = str(files[0].name).replace("/", "\\").rsplit('\\', 1)[0]
    os.chdir(p)
    pdf = Pdf.new()
    for x in files:
        path = str(x.name).replace("/", "\\")
        print(path)
        name = path.split("\\")[-1]
        with Pdf.open(path) as pd:
            pdf.pages.extend(pd.pages)

    pdf.save('merged.pdf')
    lblMergeBTN.config(text="Done!")

def unlockBT():
    lblUnlockBTN.config(text="Working!")
    files = filedialog.askopenfiles()
    p = str(files[0].name).replace("/", "\\").rsplit('\\', 1)[0]
    os.chdir(p)
    passw = simpledialog.askstring(title="Password", prompt="Enter Password")
    for x in files:
        path = str(x.name).replace("/", "\\")
        print(path)
        with Pdf.open(path, password=passw) as pd:
            pd.save(path.split()[-1].replace(".pdf", "(unlocked).pdf"))
    lblUnlockBTN.config(text="Done!")

def splitBT():
    lblSplitBTN.config(text="Working!")
    files = filedialog.askopenfiles()
    p = str(files[0].name).replace("/", "\\").rsplit('\\', 1)[0]
    os.chdir(p)
    for x in files:
        path = str(x.name).replace("/", "\\")
        print(path)
        with Pdf.open(path) as pd:
            length = len(pd.pages)
            name = path.split("\\")[-1]
            pr = "The PDF {}".format(name) + " has {} pages\n".format(length)
            promp = pr + "The splitting position. The split is inclusive the left pdf. \n So Split on 5 means 1 to 5 " \
                         "are one pdf and 6 to 10 ( pdf has 10 pages) is the second pdf. "
            intervalX = simpledialog.askstring(title="Length", prompt=promp)
            print(intervalX)
            intervalX = int(intervalX)
            pdsplit = []
            for n, page in enumerate(pd.pages):
                pdsplit.append(page)
            i = 0
            pdfsplit1 = Pdf.new()
            pdfsplit2 = Pdf.new()
            for n, x in enumerate(pdsplit):
                if i <= intervalX - 1:
                    pdfsplit1.pages.append(x)
                    i += 1
                else:
                    pdfsplit2.pages.append(x)
                    i += 1
            pdfsplit1.save(name.replace(".pdf", "(split 1 to {}).pdf".format(intervalX)))
            st = str(intervalX + 1)
            p = st + " - " + str(length)
            pdfsplit2.save(name.replace(".pdf", "(split {}).pdf".format(p)))
    lblSplitBTN.config(text="Done!")

def rotateBT():
    lblRotateBTN.config(text="Working!")
    files = filedialog.askopenfiles()
    p = str(files[0].name).replace("/", "\\").rsplit('\\', 1)[0]
    os.chdir(p)
    for x in files:
        path = str(x.name).replace("/", "\\")
        print(path)
        with Pdf.open(path) as pd:
            name = path.split("\\")[-1]
            promp = "Select degree (90/180/270) to rotate {}: ".format(name)
            rot = simpledialog.askstring(title="Rotate", prompt=promp)
            for p in pd.pages:
                p.Rotate = int(rot)
            pd.save(name.replace(".pdf", "(rotated{}).pdf".format(rot)))
    lblRotateBTN.config(text="Done!")

def cutBT():
    lblCutBTN.config(text="Working!")
    files = filedialog.askopenfiles()
    p = str(files[0].name).replace("/", "\\").rsplit('\\', 1)[0]
    os.chdir(p)
    for x in files:
        path = str(x.name).replace("/", "\\")
        print(path)
        with Pdf.open(path) as pd:
            name = path.split("\\")[-1]
            lengthPDF = len(pd.pages)
            pr = "Current PDF: {}.\n".format(name)
            prl = "Has {} pages.\n".format(lengthPDF)
            prs = "Enter an interval to be cut from pdf (From and To are inclusive).\n"
            pr0 = "If entered 0 0 the pdf will be skipped!\n"
            promp = pr + prl + prs + pr0 + "Cut From page:"
            intervalX = simpledialog.askstring(title="Cut from", prompt=promp)
            promp2 = pr + prl + "Cutting from page {} ".format(intervalX) +"Cut To page:"
            intervalY = simpledialog.askstring(title="Cut to", prompt=promp2)
            if intervalX == '0' and intervalY == '0':
                continue
            pdsplit = []
            for n, page in enumerate(pd.pages):
                pdsplit.append(page)
            i = 0
            pdfcut = Pdf.new()
            for n, f in enumerate(pdsplit):
                if (i == int(intervalX) - 1) or (i > int(intervalX) - 1 and i < int(intervalY)):
                    i += 1
                else:
                    pdfcut.pages.append(f)
                    i += 1
            intervalCut = intervalX + "-" + intervalY
            pdfcut.save(name.replace(".pdf", "(cut {}).pdf".format(intervalCut)))
    lblCutBTN.config(text="Done!")


def printBT():
    lblprintBTN.config(text="Working!")
    files = filedialog.askopenfiles()
    p = str(files[0].name).replace("/", "\\").rsplit('\\', 1)[0]
    os.chdir(p)
    for x in files:
        path = str(x.name).replace("/", "\\")
        print(path)
        with Pdf.open(path) as pd:
            length = len(pd.pages)
            name = path.split("\\")[-1]

            pdsplit  = []
            pdsplit2 = []
            for n, page in enumerate(pd.pages):
                if n % 2 == 0:
                    pdsplit.append(page)
                else:
                    pdsplit2.append(page)
            i = 0
            pdfsplit1 = Pdf.new()
            pdfsplit2 = Pdf.new()
            for n, x in enumerate(pdsplit):
               pdfsplit1.pages.append(x)

            for n, x in enumerate(pdsplit2):
               pdfsplit2.pages.append(x)

            pdfsplit1.save(name.replace(".pdf", "xSplit1.pdf"))
            pdfsplit2.save(name.replace(".pdf", "xSplit2.pdf"))
    lblprintBTN.config(text="Done!")

def exclusiveBT():
    lblexclusiveBTN.config(text="Working!")
    files = filedialog.askopenfiles()
    p = str(files[0].name).replace("/", "\\").rsplit('\\', 1)[0]
    os.chdir(p)
    for x in files:
        path = str(x.name).replace("/", "\\")
        print(path)
        with Pdf.open(path) as pd:
            name = path.split("\\")[-1]
            lengthPDF = len(pd.pages)
            pr = "Current PDF: {}.\n".format(name)
            prl = "Has {} pages.\n".format(lengthPDF)
            prs = "Enter an interval to be cut from pdf (From and To are inclusive).\n"
            pr0 = "If entered 0 0 the pdf will be skipped!\n"
            promp = pr + prl + prs + pr0 + "Cut From page:"
            intervalX = simpledialog.askstring(title="Cut from", prompt=promp)
            promp2 = pr + prl + "Cutting from page {} ".format(intervalX) +"Cut To page:"
            intervalY = simpledialog.askstring(title="Cut to", prompt=promp2)
            if intervalX == '0' and intervalY == '0':
                continue
            pdsplit = []
            for n, page in enumerate(pd.pages):
                pdsplit.append(page)
            i = 0
            pdfcut = Pdf.new()
            for n, f in enumerate(pdsplit):
                if (i != int(intervalX) - 1) and (i <= int(intervalX) - 1 or i > int(intervalY) - 1 ):
                    i += 1
                else:
                    pdfcut.pages.append(f)
                    i += 1
            intervalCut = intervalX + "-" + intervalY
            pdfcut.save(name.replace(".pdf", "(cut {}).pdf".format(intervalCut)))
    lblexclusiveBTN.config(text="Done!")

def pdf2pptBT():
    lblpdftopptBTN.config(text="Working!")
    files = filedialog.askopenfiles()
    p = str(files[0].name).replace("/", "\\").rsplit('\\', 1)[0]
    os.chdir(p)
    for x in files:
        presentation = Presentation()
        path = str(x.name).replace("/", "\\")
        print(path)
        name = path.split("\\")[-1]
        print(name)
        promp1 = "Select DPI, default 600 (over 1000 the program will freak out, recommended 300 - 800)"
        DPI = simpledialog.askstring(title="DPI", prompt=promp1)
        promp2 = "Select number of threads, default 4 (recommended 2 - 4)"
        threads = simpledialog.askstring(title="No. of Threads", prompt=promp2)
        if DPI == "" and threads != "":
            si = convert_from_path(name , 600, thread_count=int(threads))    
        elif DPI == "" and threads == "":
            si = convert_from_path(name , 600, thread_count=4)
        else:
            si = convert_from_path(name , int(DPI), thread_count=int(threads))
        
        for i, si in enumerate(si):
            imgf = BytesIO()
            si.save(imgf, format='tiff')
            data = imgf.getvalue()
            imgf.seek(0)
            width, height = si.size

            presentation.slide_height = height * 9525
            presentation.slide_width = width * 9525

            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            picture = slide.shapes.add_picture(imgf, 0, 0, width=width * 9525, height=height * 9525)
        presentation.save(name.replace(".pdf", "(pdf2ppt).ppt"))
        

    lblpdftopptBTN.config(text="Done!")


## extractor for mergechapters
def img2pdfBT():
    lblimgtopdfBTN.config(text="Working!")
    path = filedialog.askdirectory()
    mname = simpledialog.askstring(title="Name", prompt="Enter prefix name.")
    for root, dirs, files in os.walk(path):
        lst = []
        sort_nicely(files)
        for x in files:
            if x.endswith((".jpg",".png")):
                print(root +"/"+ x)
            
                image = PIL.Image.open(str(root +"/"+ x).replace("/", "\\"))
                pdf_bytes = img2pdf.convert(image.filename)
                f = io.BytesIO(pdf_bytes)
                lst.append(f)

                image.close()
            
        pdf = Pdf.new()
        for y in lst:
            with Pdf.open(y) as pd:
                pdf.pages.extend(pd.pages)
        rootname = str(root).replace("/", "\\")
        nameoffile = str(rootname + "/" + str(mname)+".pdf")
        pdf.save(nameoffile)
    lblimgtopdfBTN.config(text="Done!")


def pdf2imgBT():
    lblpdf2imgBTN.config(text="Working!")
    files = filedialog.askopenfiles()
    p = str(files[0].name).replace("/", "\\").rsplit('\\', 1)[0]
    os.chdir(p)
    filename = str(files[0].name).replace("/", "\\").rsplit('\\', 1)[1]
    images = convert_from_path(filename, int(150))
    for i in range(len(images)):
        page = 'page'+ str(i) +'.png'
        images[i].save(page, 'PNG')
        print(page)
    lblpdf2imgBTN.config(text="Done!")

import re
# https://blog.codinghorror.com/sorting-for-humans-natural-sort-order/
def sort_nicely( l ):
    convert = lambda text: int(text) if text.isdigit() else text
    alphanum_key = lambda key: [ convert(c) for c in re.split('([0-9]+)', key) ]
    return l.sort(key=alphanum_key)

def mergechaptersBT():
    lblmrgchapBTN.config(text="Working!")
    path = filedialog.askdirectory()
    os.chdir(path)
    nametolookfor = simpledialog.askstring(title="Part of string to look for", prompt="Enter substring to look for in file names. This will remove the part before the string (inclusive the string),\n eg. NAME Ch. 111 xyz.pdf -> 111 xyz.pdf\n If none type NONE and hit enter. NONE assumes correct naming for sorting order:")
    l = []
    for root, dirs, files in os.walk(path):
        for x in files:
            if nametolookfor != "NONE":
                if x.endswith("pdf") and x.find(nametolookfor) >= 0:
                    print(x.rsplit(nametolookfor)[1])
                    os.rename(x, x.rsplit(nametolookfor)[1])
                    #l.append(x)

            else:
                #l.append(x)
                next

    #sort_nicely(l)

    count = simpledialog.askinteger(title="Select count to merge", prompt="Choose how many chapters to merge:")
    name = simpledialog.askstring(title="Name of merged file", prompt="Enter the name of the merged file to be saved:")
    i = 0
    pdf = Pdf.new()
    j = 1
    #for x in l:
    for root, dirs, files in os.walk(path):
        sort_nicely(files)
        for x in files:
            if str(x).endswith(".pdf"):
                i += 1
                with open('chapters.txt', 'a') as f:
                    print(x,file=f)   
                print(x)
                with Pdf.open(x) as pd:
                    pdf.pages.extend(pd.pages)
                if i > count:
                    pdf.save(str(name+ str(j) +".pdf"))
                    pdf = Pdf.new()
                    i = 0
                    j+= 1
                    with open('chapters.txt', 'a') as f:
                        print(file=f)
                        print(j,file=f)
                        print("-----------------------------------------------------------------------------",file=f)
                        print(file=f)
                    print()
                    print(j)
                    print("-----------------------------------------------------------------------------")
                    print()
    if i != count:
        pdf.save(str(name+ str(j) +".pdf"))
    lblmrgchapBTN.config(text="Done!")

def reducePDFsize():
    lblreduceBTN.config(text="Working!")
    files = filedialog.askopenfiles()
    p = str(files[0].name).replace("/", "\\").rsplit('\\', 1)[0]
    os.chdir(p)
    filename = str(files[0].name).replace("/", "\\").rsplit('\\', 1)[1]
    images = convert_from_path(filename, int(150))
    pngs = []
    for i in range(len(images)):
        page = 'page'+ str(i) +'.png'
        images[i].save(page, 'PNG')
        print(page)
        pngs.append(page)
    lst = []
    for root, dirs, files in os.walk(p):
        sort_nicely(files)
        for x in files:
            if x.endswith((".jpg",".png")):
                image = PIL.Image.open(str(root +"/"+ x).replace("/", "\\"))
                pdf_bytes = img2pdf.convert(image.filename)
                f = io.BytesIO(pdf_bytes)
                lst.append(f)

                image.close()
            
    pdf = Pdf.new()
    for y in lst:
        with Pdf.open(y) as pd:
            pdf.pages.extend(pd.pages)

    nameoffile = str("reduced" +".pdf")
    pdf.save(nameoffile)
    for x in pngs:
        os.remove(x)
    lblreduceBTN.config(text="Done!")


if __name__ == '__main__':
    root = Tk()
    root.geometry('290x550')
    root.title('PDF Mastery')
    root['bg'] = '#000000'
    root.resizable(False, False)

    mergeBTN = Button(root, text="Merge PDFs", command=mergeBT)
    mergeBTN.grid(row=9, column=0, stick=W, pady=10, padx=50)

    unlockBTN = Button(root, text="Unlock PDFs", command=unlockBT)
    unlockBTN.grid(row=19, column=0, stick=W, pady=10, padx=50)

    splitBTN = Button(root, text="Split PDFs", command=splitBT)
    splitBTN.grid(row=29, column=0, stick=W, pady=10, padx=50)

    rotateBTN = Button(root, text="Rotate PDFs", command=rotateBT)
    rotateBTN.grid(row=39, column=0, stick=W, pady=10, padx=50)

    cutBTN = Button(root, text="Cut PDFs", command=cutBT)
    cutBTN.grid(row=49, column=0, stick=W, pady=10, padx=50)

    printBTN = Button(root, text="Print Split", command=printBT)
    printBTN.grid(row=59, column=0, stick=W, pady=10, padx=50)

    exclusiveCUTBTN = Button(root, text="Exclusive Cut", command=exclusiveBT)
    exclusiveCUTBTN.grid(row=69, column=0, stick=W, pady=10, padx=50)

    pdftopptxBTN = Button(root, text="PDF -> PPT", command=pdf2pptBT)
    pdftopptxBTN.grid(row=79, column=0, stick=W, pady=10, padx=50)

    imgtopdfBTN = Button(root, text="Img -> PDF", command=img2pdfBT)
    imgtopdfBTN.grid(row=89, column=0, stick=W, pady=10, padx=50)

    mrgchapBTN = Button(root, text="Bulk mrg chapt", command=mergechaptersBT)
    mrgchapBTN.grid(row=99, column=0, stick=W, pady=10, padx=50)

    pdf2imgBTN = Button(root, text="pdf -> img", command=pdf2imgBT)
    pdf2imgBTN.grid(row=109, column=0, stick=W, pady=10, padx=50)
    
    reduceBTN = Button(root, text="reduce size", command=reducePDFsize)
    reduceBTN.grid(row=119, column=0, stick=W, pady=10, padx=50)

    lblMergeBTN = Label(root, text='', width=10)
    lblMergeBTN.grid(row=9, column=1)

    lblCutBTN = Label(root, text='', width=10)
    lblCutBTN.grid(row=49, column=1)

    lblSplitBTN = Label(root, text='', width=10)
    lblSplitBTN.grid(row=29, column=1)

    lblRotateBTN = Label(root, text='', width=10)
    lblRotateBTN.grid(row=39, column=1)

    lblUnlockBTN = Label(root, text='', width=10)
    lblUnlockBTN.grid(row=19, column=1)

    lblprintBTN = Label(root, text='', width=10)
    lblprintBTN.grid(row=59, column=1)

    lblexclusiveBTN = Label(root, text='', width=10)
    lblexclusiveBTN.grid(row=69, column=1)

    lblpdftopptBTN = Label(root, text='', width=10)
    lblpdftopptBTN.grid(row=79, column=1)

    lblimgtopdfBTN = Label(root, text='', width=10)
    lblimgtopdfBTN.grid(row=89, column=1)

    lblmrgchapBTN = Label(root, text='', width=10)
    lblmrgchapBTN.grid(row=99, column=1)

    lblpdf2imgBTN = Label(root, text='', width=10)
    lblpdf2imgBTN.grid(row=109, column=1)

    lblreduceBTN = Label(root, text='', width=10)
    lblreduceBTN.grid(row=119, column=1)

    root.mainloop()

   
